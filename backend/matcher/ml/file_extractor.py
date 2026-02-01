"""
File extraction utilities for handling different file types
Supports: PDF, DOCX, TXT, DOC
"""

import os
import re
import tempfile
from typing import Optional, Tuple

try:
    import PyPDF2
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import python_docx
    HAS_PYTHON_DOCX = True
except ImportError:
    HAS_PYTHON_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class FileExtractionError(Exception):
    """Custom exception for file extraction errors"""
    pass


class FileExtractor:
    """Extract text from various file formats"""

    SUPPORTED_FORMATS = {
        'pdf': '.pdf',
        'docx': '.docx',
        'doc': '.doc',
        'txt': '.txt',
        'text': '.txt',
    }

    MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB

    @staticmethod
    def extract_pdf(file_path: str) -> str:
        """Extract text from PDF file"""
        if not HAS_PYPDF:
            raise FileExtractionError(
                "PyPDF2 not installed. Install with: pip install PyPDF2"
            )

        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text()
            return text.strip()
        except Exception as e:
            raise FileExtractionError(f"Error extracting PDF: {str(e)}")

    @staticmethod
    def extract_docx(file_path: str) -> str:
        """Extract text from DOCX file"""
        if not HAS_DOCX:
            raise FileExtractionError(
                "python-docx not installed. Install with: pip install python-docx"
            )

        try:
            doc = Document(file_path)
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            # Extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            
            return text.strip()
        except Exception as e:
            raise FileExtractionError(f"Error extracting DOCX: {str(e)}")

    @staticmethod
    def extract_doc(file_path: str) -> str:
        """Extract text from DOC file using python-docx (modern .doc files)"""
        # Modern .doc files are often handled by python-docx
        return FileExtractor.extract_docx(file_path)

    @staticmethod
    def extract_txt(file_path: str) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read().strip()
        except Exception as e:
            raise FileExtractionError(f"Error extracting TXT: {str(e)}")

    @staticmethod
    def extract_pptx(file_path: str) -> str:
        """Extract text from PPTX file"""
        if not HAS_PPTX:
            raise FileExtractionError(
                "python-pptx not installed. Install with: pip install python-pptx"
            )

        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text.strip()
        except Exception as e:
            raise FileExtractionError(f"Error extracting PPTX: {str(e)}")

    @classmethod
    def validate_file(cls, file) -> Tuple[bool, str]:
        """
        Validate file before extraction
        Returns: (is_valid, error_message)
        """
        if not file:
            return False, "No file provided"

        # Check file size
        if file.size > cls.MAX_FILE_SIZE:
            return False, f"File size exceeds {cls.MAX_FILE_SIZE / (1024*1024):.0f} MB limit"

        # Get file extension
        file_name = file.name.lower()
        file_ext = os.path.splitext(file_name)[1]

        # Check if format is supported
        supported_exts = ['.pdf', '.docx', '.doc', '.txt', '.pptx']
        if file_ext not in supported_exts:
            return False, f"Unsupported file format. Supported: {', '.join(supported_exts)}"

        return True, ""

    @classmethod
    def extract_text(cls, file) -> str:
        """
        Main method to extract text from file
        Automatically detects file type
        """
        # Validate file
        is_valid, error_msg = cls.validate_file(file)
        if not is_valid:
            raise FileExtractionError(error_msg)

        # Save file to system temporary directory (cross-platform)
        temp_dir = tempfile.gettempdir()
        temp_path = os.path.join(temp_dir, file.name)
        try:
            # Save uploaded file
            with open(temp_path, 'wb') as temp_file:
                for chunk in file.chunks():
                    temp_file.write(chunk)

            # Extract based on file type
            file_ext = os.path.splitext(file.name.lower())[1]

            if file_ext == '.pdf':
                text = cls.extract_pdf(temp_path)
            elif file_ext == '.docx':
                text = cls.extract_docx(temp_path)
            elif file_ext == '.doc':
                text = cls.extract_doc(temp_path)
            elif file_ext == '.txt':
                text = cls.extract_txt(temp_path)
            elif file_ext == '.pptx':
                text = cls.extract_pptx(temp_path)
            else:
                raise FileExtractionError(f"Unknown file format: {file_ext}")

            if not text or len(text.strip()) < 10:
                raise FileExtractionError(
                    "Extracted text is too short or empty. "
                    "Please provide a document with meaningful content."
                )

            return text

        finally:
            # Clean up temporary file
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except:
                    pass
